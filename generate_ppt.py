#!/usr/bin/env python3

from fastmcp import FastMCP
from pydantic import Field, BaseModel
from pptx import Presentation
import os

mcp = FastMCP("ptms made easier")

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "template.pptx")
OUTPUT_FOLDER = os.path.join(SCRIPT_DIR, "generated_ppts")
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


class StudentData(BaseModel):
    Name: str = Field(..., description="Student's name")
    Class: str = Field(..., description="Class")
    Subjects: str = Field(..., description="Student's subject")
    Grade: str = Field(..., description="Student's grade")
    Comments: str = Field(..., description="comments you might want to leave")


def add_stuff_to_ppt(prs: Presentation, data: StudentData):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in data.dict().items():
                    if f"{{{{{key}}}}}" in shape.text:
                        if key == "Subjects":
                            value = ", ".join([s.strip() for s in value.split(",")])
                        shape.text = shape.text.replace(f"{{{{{key}}}}}", str(value))


@mcp.tool
def generate_ppt(student: StudentData):
    """Genrate a ppt for your lovely students based on the template"""
    prs = Presentation(TEMPLATE_PATH)
    add_stuff_to_ppt(prs, student)
    output_path = os.path.join(OUTPUT_FOLDER, f"{student.Name}.pptx")
    prs.save(output_path)
    return f"PPT generated for {student.Name} at {output_path}"



def main():
    mcp.run()

if __name__=="__main__":
    main()

